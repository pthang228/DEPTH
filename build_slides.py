# -*- coding: utf-8 -*-
"""DEPTH deck v3 — deep theory + detailed diagrams + one running numeric example."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

NAVY  = RGBColor(0x14,0x21,0x3D); NAVY2 = RGBColor(0x1E,0x2F,0x55)
LIGHT = RGBColor(0xF4,0xF6,0xFB); WHITE = RGBColor(0xFF,0xFF,0xFF)
INK   = RGBColor(0x20,0x2C,0x44); MUTE  = RGBColor(0x6B,0x77,0x93)
ICE   = RGBColor(0xCA,0xDC,0xFC)
BLUE  = RGBColor(0x3B,0x6F,0xD4); TEAL  = RGBColor(0x12,0xA8,0x9D); GOLD = RGBColor(0xE8,0xA3,0x3D)
SOFTBLUE=RGBColor(0xE7,0xEE,0xFB); SOFTTEAL=RGBColor(0xE1,0xF4,0xF2); SOFTGOLD=RGBColor(0xFB,0xF1,0xDF)
CODEBG = RGBColor(0x0E,0x1A,0x33); LINEC = RGBColor(0xD7,0xDE,0xEC); ARROWC=RGBColor(0xAA,0xB4,0xC8)
HF,BF,CF = "Georgia","Calibri","Consolas"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)
def bg(s,c): s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def Rn(t,sz,c,b=False,it=False,f=BF): return (t,sz,c,b,it,f)
def shp(s,k,x,y,w,h,fill,line=None,lw=1.0):
    sp=s.shapes.add_shape(k,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def rect(s,x,y,w,h,fill,line=None,lw=1.0): return shp(s,MSO_SHAPE.RECTANGLE,x,y,w,h,fill,line,lw)
def rrect(s,x,y,w,h,fill,line=None,lw=1.0): return shp(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h,fill,line,lw)
def txt(s,x,y,w,h,paras,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    for i,p in enumerate(paras):
        pa=tf.paragraphs[0] if i==0 else tf.add_paragraph(); pa.alignment=align
        if isinstance(p,dict):
            if p.get('space_after') is not None: pa.space_after=Pt(p['space_after'])
            if p.get('space_before') is not None: pa.space_before=Pt(p['space_before'])
            if p.get('line') is not None: pa.line_spacing=p['line']
            p=p['runs']
        for (t,sz,c,b,it,f) in p:
            r=pa.add_run(); r.text=t; r.font.size=Pt(sz); r.font.color.rgb=c
            r.font.bold=b; r.font.italic=it; r.font.name=f
    return tb
def header(s,eyebrow,title,accent=GOLD):
    rect(s,0,0,13.333,1.32,LIGHT); rect(s,0.6,0.34,0.12,0.6,accent)
    txt(s,0.85,0.32,11.8,0.3,[[Rn(eyebrow.upper(),11.5,accent,True)]])
    txt(s,0.85,0.58,12.1,0.56,[[Rn(title,25,INK,True,f=HF)]])
def chip(s,x,y,w,h,color,label,tcol=WHITE,size=13):
    rrect(s,x,y,w,h,color); txt(s,x,y,w,h,[[Rn(label,size,tcol,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def badge(s,x,y,d,color,n):
    shp(s,MSO_SHAPE.OVAL,x,y,d,d,color); txt(s,x,y,d,d,[[Rn(n,16,WHITE,True,f=HF)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def arrow(s,x,y,w,h,c=None): shp(s,MSO_SHAPE.RIGHT_ARROW,x,y,w,h,c or ARROWC)
def darrow(s,x,y,w,h,c=None): shp(s,MSO_SHAPE.DOWN_ARROW,x,y,w,h,c or ARROWC)
def connect(s,x1,y1,x2,y2,color,wd=2.0):
    cn=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    cn.line.color.rgb=color; cn.line.width=Pt(wd); cn.shadow.inherit=False; return cn
def formula(s,x,y,w,h,lines,title=None):
    rrect(s,x,y,w,h,CODEBG,CODEBG)
    paras=[]
    if title: paras.append({'runs':[Rn(title,11.5,GOLD,True,f=BF)],'space_after':6})
    for ln in lines:
        paras.append({'runs':[Rn(ln,13,ICE,f=CF)],'space_after':5})
    txt(s,x+0.2,y+0.14,w-0.4,h-0.24,paras,anchor=MSO_ANCHOR.MIDDLE)
def card(s,x,y,w,h,fill=WHITE,line=LINEC,lw=1): return rrect(s,x,y,w,h,fill,line,lw)
def bullets(s,x,y,w,h,items,col=INK,bcol=None,size=13.5,gap=6):
    bcol=bcol or GOLD
    paras=[]
    for it in items:
        paras.append({'runs':[Rn("•  ",size,bcol,True),Rn(it,size,col)],'space_after':gap})
    txt(s,x,y,w,h,paras)
def divider(part,title,sub,accent):
    s=slide(); bg(s,NAVY)
    shp(s,MSO_SHAPE.OVAL,0.9,2.45,0.95,0.95,accent)
    txt(s,0.9,2.45,0.95,0.95,[[Rn(part,26,NAVY if accent==GOLD else WHITE,True,f=HF)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,2.15,2.4,10.4,0.9,[[Rn(title,38,WHITE,True,f=HF)]])
    rect(s,2.2,3.4,3.0,0.05,accent)
    txt(s,2.18,3.58,10.4,0.7,[[Rn(sub,17,ICE,it=True)]])
    return s
# example panel helper (right side, gold)
def example_panel(s,x,y,w,h,title,paras):
    rrect(s,x,y,w,h,SOFTGOLD,GOLD,1.25)
    txt(s,x+0.22,y+0.16,w-0.44,0.4,[[Rn("VÍ DỤ XUYÊN SUỐT",11,GOLD,True),Rn("  "+title,11.5,INK,True)]])
    txt(s,x+0.22,y+0.62,w-0.44,h-0.78,paras)

NOTES=[]
def note(t): NOTES.append(t)

# ===================================================== 1 TITLE
s=slide(); bg(s,NAVY)
rect(s,0,0,13.333,0.18,GOLD); rect(s,0,7.32,13.333,0.18,TEAL)
txt(s,0.9,1.4,11.6,0.4,[[Rn("TEMPORAL KNOWLEDGE GRAPH EXTRAPOLATION  ·  ICEWS14",13,ICE,True)]])
txt(s,0.85,1.95,11.6,1.5,[[Rn("DEPTH",92,WHITE,True,f=HF)]])
txt(s,0.9,3.45,11.6,0.6,[[Rn("Dual Evolution–Path Temporal Reasoning",22,GOLD,True,it=True,f=HF)]])
txt(s,0.9,4.25,11.4,1.0,[[Rn("Hợp nhất ",16,ICE),Rn("suy luận theo tiến hóa nhúng",16,WHITE,True),
  Rn(" và ",16,ICE),Rn("suy luận theo đường quan hệ",16,WHITE,True),Rn(" — huấn luyện end-to-end.",16,ICE)]])
chip(s,0.9,5.55,2.5,0.5,BLUE,"Nhánh Nhúng"); chip(s,3.55,5.55,2.5,0.5,TEAL,"Nhánh Đường"); chip(s,6.2,5.55,2.5,0.5,GOLD,"Hợp nhất",tcol=NAVY)
note("Bìa. DEPTH gộp 2 lối suy luận trên TKG; deck đi sâu lý thuyết + 1 ví dụ số chạy xuyên suốt.")

# ===================================================== 2 OUTLINE
s=slide(); bg(s,LIGHT)
txt(s,0.85,0.5,11,0.8,[[Rn("Nội dung",30,INK,True,f=HF)]])
items=[("1","Nền tảng & ví dụ","Bài toán, ký hiệu, và ví dụ số dùng xuyên suốt",NAVY2,ICE),
       ("2","LogCL — nhánh nhúng","RGCN · GRU · attention · contrastive · ConvTransE",BLUE,SOFTBLUE),
       ("3","DEPTH — hợp nhất","Backbone chung · Head A/B · path reasoner · fusion γ",GOLD,SOFTGOLD),
       ("4","Vì sao tốt hơn & Hiệu năng","Lý do + bảng kết quả ICEWS14",TEAL,SOFTTEAL)]
y=1.6
for n,t,d,c,soft in items:
    rrect(s,0.85,y,11.6,1.08,soft); badge(s,1.15,y+0.27,0.55,c,n)
    txt(s,2.05,y+0.18,9.9,0.45,[[Rn(t,19,INK,True,f=HF)]]); txt(s,2.05,y+0.62,10.0,0.4,[[Rn(d,13.5,MUTE)]])
    y+=1.3
note("Mục lục. 4 phần; nhấn rằng có 1 ví dụ số theo suốt deck.")

# ===================================================== 3 Bài toán & ký hiệu
s=slide(); bg(s,LIGHT); header(s,"Phần 1 · Nền tảng","Bài toán & ký hiệu",NAVY2)
txt(s,0.85,1.55,7.1,4.9,[
 {'runs':[Rn("TKG = chuỗi snapshot theo thời gian:  ",14.5,INK),Rn("G₁, G₂, …, G_T",14.5,NAVY2,True)],'space_after':7},
 {'runs':[Rn("Mỗi G_t = tập sự kiện (s, r, o) xảy ra tại thời điểm t.",14,INK)],'space_after':12},
 {'runs':[Rn("Nhiệm vụ (extrapolation): ",14.5,INK,True),
          Rn("cho truy vấn (s, r, ?, t) ở thời điểm TƯƠNG LAI t, xếp hạng mọi thực thể làm object.",14,INK)],'space_after':12},
 {'runs':[Rn("Đánh giá: ",14.5,INK,True),Rn("MRR, Hits@1/3/10, theo time-aware filtered.",14,INK)],'space_after':6},
])
formula(s,0.85,5.0,7.1,1.45,[
 "N = 7128 thực thể,  R quan hệ (+nghịch đảo → 2R)",
 "d = 200 (chiều embedding),  K = 7 (snapshot lịch sử)",
 "H ∈ ℝ^(N×d): embedding thực thể;  emb_rel ∈ ℝ^(2R×d)"],title="Ký hiệu dùng xuyên suốt")
rrect(s,8.25,1.55,4.25,4.9,NAVY,NAVY)
txt(s,8.55,1.8,3.7,0.5,[[Rn("Hai họ phương pháp",15,GOLD,True,f=HF)]])
txt(s,8.55,2.35,3.7,4.0,[
 {'runs':[Rn("① Tiến hóa nhúng",14,ICE,True)],'space_after':2},
 {'runs':[Rn("học vector & cho tiến hóa theo thời gian (LogCL, RE-GCN).",12.5,ICE)],'space_after':12},
 {'runs':[Rn("② Suy luận đường",14,ICE,True)],'space_after':2},
 {'runs':[Rn("truy vết đường quan hệ đa-hop trong lịch sử (CognTKE).",12.5,ICE)],'space_after':12},
 {'runs':[Rn("DEPTH = gộp cả hai.",13.5,GOLD,True)]},
])
note("Định nghĩa bài toán & ký hiệu. Nói rõ extrapolation (dự đoán tương lai), metric, kích thước N,d,K. Giới thiệu 2 họ phương pháp mà DEPTH gộp.")

# ===================================================== 4 Ví dụ xuyên suốt
s=slide(); bg(s,LIGHT); header(s,"Phần 1 · Nền tảng","Ví dụ minh họa cho phần hợp nhất (DEPTH)",GOLD)
rrect(s,0.85,1.5,5.6,4.9,SOFTGOLD,GOLD,1.25)
txt(s,1.1,1.7,5.1,0.5,[[Rn("Truy vấn",13,GOLD,True),Rn("   (USA, Sanction, ?, t)",16,INK,True,f=HF)]])
txt(s,1.1,2.25,5.1,0.4,[[Rn("Đáp án đúng: ",13.5,INK),Rn("Iran",15,TEAL,True)]])
txt(s,1.1,2.75,5.1,0.4,[[Rn("4 ứng viên ta theo dõi:",13.5,INK,True)]])
for i,(e,c) in enumerate([("Iran (đúng)",TEAL),("Russia",INK),("China",INK),("North Korea",INK)]):
    txt(s,1.3,3.2+i*0.45,5.0,0.4,[[Rn("•  ",13,GOLD,True),Rn(e,13.5,c,e=="Iran (đúng)")]] if False else [[Rn("•  ",13,GOLD,True),Rn(e,13.5,c,c==TEAL)]])
txt(s,1.1,5.2,5.1,1.1,[[Rn("Lịch sử gần có các sự kiện: USA→threaten→Iran, Iran→use_force, …  (không có đường rõ tới Russia/China).",12.5,MUTE,it=True)]])
rrect(s,6.75,1.5,5.7,4.9,NAVY,NAVY)
txt(s,7.05,1.72,5.1,0.5,[[Rn("Ta sẽ theo ví dụ này qua từng bước:",14,GOLD,True,f=HF)]])
steps=[("LogCL → S_embed","embedding chấm điểm: Russia cao nhất (sai!), Iran chỉ #3"),
       ("Head B → S_path","đường quan hệ tìm thấy Iran → điểm path cao"),
       ("Fusion → kết quả","cộng lại: Iran vươn lên #1 (đúng)")]
for i,(t,d) in enumerate(steps):
    yy=2.45+i*1.25
    badge(s,7.05,yy,0.5,GOLD,str(i+1))
    txt(s,7.75,yy-0.05,4.4,0.4,[[Rn(t,14.5,WHITE,True,f=HF)]])
    txt(s,7.75,yy+0.38,4.45,0.7,[[Rn(d,12.5,ICE)]])
note("Giới thiệu ví dụ số: (USA, Sanction, ?, Iran). 4 ứng viên. Báo trước câu chuyện: embedding đoán nhầm Russia, path tìm ra Iran, fusion sửa lại. Ví dụ này quay lại ở slide 12, 18, 20.")

# ===================================================== 5 DIVIDER P2 LogCL
divider("2","LogCL","Logic-aware Contrastive learning (ICDE 2024): 2 vấn đề · 3 thành phần",BLUE)
note("Chuyển sang LogCL: trình bày 2 vấn đề paper nhắm tới, 3 thành phần kiến trúc, ví dụ end-to-end và kết quả.")

# màu cho sơ đồ (theo SVG)
PURPLE=RGBColor(0x53,0x4A,0xB7); SOFTPUR=RGBColor(0xEE,0xED,0xFE)
RUST=RGBColor(0x99,0x3C,0x1D);  SOFTRUST=RGBColor(0xFA,0xEC,0xE7)
GREEN=RGBColor(0x0F,0x6E,0x56); SOFTGRN=RGBColor(0xE1,0xF5,0xEE)

# ===================================================== 6 Hai vấn đề
s=slide(); bg(s,LIGHT); header(s,"Phần 2 · LogCL","Hai vấn đề LogCL nhắm tới",BLUE)
rrect(s,0.85,1.75,5.65,4.6,SOFTBLUE,SOFTBLUE); rect(s,0.85,1.75,5.65,0.1,BLUE)
txt(s,1.1,1.95,5.15,0.5,[[Rn("Vấn đề 1",12,BLUE,True),Rn("  Lịch sử liên quan bị bỏ qua",15,INK,True,f=HF)]])
txt(s,1.1,2.6,5.2,3.5,[
 {'runs':[Rn("Model cũ mặc định snapshot GẦN nhất là quan trọng nhất.",14,INK)],'space_after':10},
 {'runs':[Rn("Nhưng entity trong query (vd China) có thể KHÔNG xuất hiện ở snapshot gần → snapshot đó vô ích.",14,INK)],'space_after':10},
 {'runs':[Rn("Trong khi một snapshot CŨ hơn lại rất liên quan → bị bỏ lỡ.",14,INK)]},
])
rrect(s,6.8,1.75,5.65,4.6,SOFTRUST,SOFTRUST); rect(s,6.8,1.75,5.65,0.1,RUST)
txt(s,7.05,1.95,5.15,0.5,[[Rn("Vấn đề 2",12,RUST,True),Rn("  Yếu trước nhiễu",15,INK,True,f=HF)]])
txt(s,7.05,2.6,5.2,1.6,[
 {'runs':[Rn("Thêm Gaussian noise vào đầu vào → MRR của RE-GCN sụt giảm tới ",14,INK),
          Rn("63.8%",15,RUST,True),Rn(".",14,INK)],'space_after':8},
 {'runs':[Rn("Mô hình quá dễ bị “lung lay” bởi chi tiết bề mặt.",14,INK)]},
])
rrect(s,7.05,4.35,5.15,1.7,WHITE,RUST,1)
txt(s,7.3,4.55,4.7,1.3,[
 {'runs':[Rn("→ LogCL giải quyết bằng:",13.5,RUST,True)],'space_after':6},
 {'runs':[Rn("• Entity-Aware Attention (chọn đúng snapshot liên quan)",13,INK)],'space_after':4},
 {'runs':[Rn("• Local–Global Contrastive (bền với noise)",13,INK)]},
])
note("Hai vấn đề: (1) lịch sử liên quan tới query bị bỏ qua vì model cũ chỉ trọng snapshot gần; (2) yếu trước noise — RE-GCN sụt 63.8% MRR khi thêm Gaussian noise. LogCL khắc phục bằng entity-aware attention + contrastive.")

# ===================================================== 7 SƠ ĐỒ kiến trúc (vẽ lại SVG)
s=slide(); bg(s,LIGHT); header(s,"Phần 2 · LogCL","Sơ đồ kiến trúc tổng quan",BLUE)
# input bar
rrect(s,1.4,1.45,10.5,0.6,RGBColor(0xF1,0xEF,0xE8),RGBColor(0x5F,0x5E,0x5A),0.75)
txt(s,1.4,1.45,10.5,0.6,[[Rn("Temporal Knowledge Graph — chuỗi snapshot (e_s, r, e_o, t)",13,INK,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
darrow(s,3.4,2.05,0.34,0.32); darrow(s,9.6,2.05,0.34,0.32)
txt(s,2.0,2.05,1.7,0.3,[[Rn("m snapshot gần",10.5,MUTE)]],align=PP_ALIGN.CENTER)
txt(s,9.7,2.05,2.0,0.3,[[Rn("toàn bộ lịch sử",10.5,MUTE)]],align=PP_ALIGN.CENTER)
# Local encoder box
rrect(s,0.7,2.45,5.7,2.0,SOFTGRN,GREEN,1.0)
txt(s,0.9,2.55,5.3,0.4,[[Rn("Local Entity-Aware Attention Recurrent Encoder",12.5,GREEN,True,f=HF)]],align=PP_ALIGN.CENTER)
rrect(s,0.95,3.05,2.55,0.55,WHITE,GREEN,0.75); txt(s,0.95,3.05,2.55,0.55,[[Rn("① R-GCN Local",11.5,GREEN,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
arrow(s,3.52,3.18,0.3,0.28,GREEN)
rrect(s,3.85,3.05,2.4,0.55,WHITE,GREEN,0.75); txt(s,3.85,3.05,2.4,0.55,[[Rn("② GRU Ent",11.5,GREEN,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
darrow(s,3.4,3.62,0.28,0.22,GREEN)
rrect(s,0.95,3.78,5.3,0.55,WHITE,GREEN,0.75); txt(s,0.95,3.78,5.3,0.55,[[Rn("③ Entity-Aware Attention  (αᵢ theo snapshot)",11.5,GREEN,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# Global encoder box
rrect(s,6.9,2.45,5.7,2.0,SOFTPUR,PURPLE,1.0)
txt(s,7.1,2.55,5.3,0.4,[[Rn("Global Entity-Aware Attention Encoder",12.5,PURPLE,True,f=HF)]],align=PP_ALIGN.CENTER)
rrect(s,7.15,3.05,5.2,0.55,WHITE,PURPLE,0.75); txt(s,7.15,3.05,5.2,0.55,[[Rn("① Xây subgraph lịch sử (1-hop + 2-hop)",11.5,PURPLE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
rrect(s,7.15,3.78,2.5,0.55,WHITE,PURPLE,0.75); txt(s,7.15,3.78,2.5,0.55,[[Rn("② R-GCN Global",11.5,PURPLE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
arrow(s,9.67,3.9,0.3,0.28,PURPLE)
rrect(s,9.85,3.78,2.5,0.55,WHITE,PURPLE,0.75); txt(s,9.85,3.78,2.5,0.55,[[Rn("③ Entity Attention",11.5,PURPLE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# contrastive
darrow(s,3.4,4.47,0.3,0.3,GREEN); darrow(s,9.6,4.47,0.3,0.3,PURPLE)
rrect(s,3.9,4.85,5.5,0.62,SOFTRUST,RUST,1.0)
txt(s,3.9,4.85,5.5,0.62,[[Rn("Local–Global Query Contrast Module",13,RUST,True,f=HF)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
connect(s,3.6,4.6,5.3,4.85,RUST,1.5); connect(s,9.5,4.6,7.9,4.85,RUST,1.5)
# loss row
rrect(s,3.0,5.7,2.9,0.55,RGBColor(0xFA,0xEE,0xDA),RGBColor(0x85,0x4F,0x0B),0.75)
txt(s,3.0,5.7,2.9,0.55,[[Rn("L_cl  (chống noise)",11.5,RGBColor(0x85,0x4F,0x0B),True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
rrect(s,7.45,5.7,2.9,0.55,RGBColor(0xE6,0xF1,0xFB),RGBColor(0x18,0x5F,0xA5),0.75)
txt(s,7.45,5.7,2.9,0.55,[[Rn("L_tkg  (dự đoán)",11.5,RGBColor(0x18,0x5F,0xA5),True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
connect(s,5.4,5.47,4.6,5.7,MUTE,1.5); connect(s,7.9,5.47,8.7,5.7,MUTE,1.5)
rrect(s,5.5,6.5,2.3,0.5,RGBColor(0xEA,0xF3,0xDE),RGBColor(0x3B,0x6D,0x11),0.75)
txt(s,5.5,6.5,2.3,0.5,[[Rn("L = L_tkg + L_cl",12.5,RGBColor(0x3B,0x6D,0x11),True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
connect(s,4.45,6.25,6.0,6.5,MUTE,1.5); connect(s,8.9,6.25,7.3,6.5,MUTE,1.5)
note("Sơ đồ (vẽ lại theo SVG): TKG → Local Encoder (R-GCN Local → GRU → Entity-Aware Attention) và Global Encoder (subgraph 1+2hop → R-GCN Global → Entity Attention); cả hai vào Contrast Module; sinh L_cl (chống noise) + L_tkg (dự đoán) → L = L_tkg + L_cl.")

# ===================================================== 8 Thành phần 1
s=slide(); bg(s,LIGHT); header(s,"Phần 2 · LogCL · Thành phần 1","Local Entity-Aware Attention Recurrent Encoder",GREEN if False else BLUE)
txt(s,0.85,1.55,7.1,0.45,[[Rn("Xử lý m snapshot gần nhất (thường 7–9 bước).",13.5,MUTE,it=True)]])
steps=[("① R-GCN Local","Lan truyền thông tin giữa node theo cạnh quan hệ TRONG mỗi snapshot. Kèm mã hóa thời gian cos(d·w + b) để biết snapshot cách query bao lâu."),
       ("② GRU","Xâu chuỗi các snapshot theo thứ tự thời gian — như đọc lịch sử từng trang."),
       ("③ Entity-Aware Attention (mới)","Tính điểm αᵢ cho TỪNG snapshot, dựa trên entity của query có xuất hiện/liên quan tới snapshot đó không. Snapshot liên quan → “chú ý” nhiều hơn.")]
y=2.1
for i,(t,d) in enumerate(steps):
    rrect(s,0.85,y,7.1,1.25,WHITE,LINEC,1); rect(s,0.85,y,0.1,1.25,BLUE)
    txt(s,1.1,y+0.13,6.7,0.4,[[Rn(t,14.5,INK,True,f=HF)]])
    txt(s,1.1,y+0.55,6.75,0.65,[[Rn(d,12.5,MUTE)]])
    y+=1.4
formula(s,8.15,1.6,4.35,1.7,[
 "time-enc:  φ(d) = cos(d·w + b)",
 "αᵢ = softmax( score(query, snapshotᵢ) )",
 "h_local = Σᵢ αᵢ · hᵢ"],title="Công thức chính")
example_panel(s,8.15,3.5,4.35,2.85,"αᵢ (China, Cooperate)",[
 {'runs':[Rn("t_q−2: China VẮNG → αᵢ thấp (gần như bỏ qua).",12.5,INK)],'space_after':8},
 {'runs':[Rn("t_q−1: China “Make a visit” → αᵢ cao (chú ý mạnh).",12.5,INK)],'space_after':8},
 {'runs':[Rn("Đây là đóng góp lớn nhất theo ablation.",12.5,GOLD,True)]},
])
note("Thành phần 1 (Local): R-GCN local + time-enc cos(dw+b); GRU nối snapshot; Entity-Aware Attention tính αᵢ theo việc entity query có liên quan snapshot — phần MỚI, đóng góp lớn nhất (ablation -8~10% nếu bỏ).")

# ===================================================== 9 Thành phần 2
s=slide(); bg(s,LIGHT); header(s,"Phần 2 · LogCL · Thành phần 2","Global Entity-Aware Attention Encoder",PURPLE)
txt(s,0.85,1.55,7.1,0.45,[[Rn("Nhìn TOÀN BỘ lịch sử, không chỉ m bước gần đây.",13.5,MUTE,it=True)]])
bullets(s,0.85,2.1,7.1,3.0,[
 "Xây subgraph lịch sử: gom mọi sự kiện liên quan entity của query — cả 1-hop lẫn 2-hop (bạn của bạn).",
 "R-GCN Global xử lý subgraph này.",
 "Một lớp attention chọn phần quan trọng nhất với query.",
 "Bắt được pattern lặp lại từ rất xa trong quá khứ."],col=INK,bcol=PURPLE,size=13.5,gap=10)
example_panel(s,8.15,1.6,4.35,4.75,"China + Cooperate",[
 {'runs':[Rn("Subgraph chứa các đối tác hợp tác LÂU DÀI của Trung Quốc.",12.5,INK)],'space_after':10},
 {'runs':[Rn("Attention chọn entity xuất hiện nhiều lần trong pattern “Cooperate”.",12.5,INK)],'space_after':10},
 {'runs':[Rn("Nếu TQ thường “họp bàn” trước khi “ký thỏa thuận”, các sự kiện chuẩn bị từ lâu vẫn được bắt → h_g.",12.5,INK)]},
])
note("Thành phần 2 (Global): xây subgraph lịch sử 1-hop+2-hop quanh entity query, R-GCN global + attention chọn phần liên quan. Bắt pattern lặp lại từ xa. Ví dụ China: đối tác hợp tác lâu dài, sự kiện 'họp bàn' chuẩn bị.")

# ===================================================== 10 Thành phần 3
s=slide(); bg(s,LIGHT); header(s,"Phần 2 · LogCL · Thành phần 3","Local–Global Query Contrast (chống noise)",RUST)
# pull/push mini
def onodeC(cx,cy,d,fill,line,label,tc,sz=12):
    shp(s,MSO_SHAPE.OVAL,cx-d/2,cy-d/2,d,d,fill,line,1.5)
    txt(s,cx-d/2,cy-d/2,d,d,[[Rn(label,sz,tc,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
onodeC(2.0,2.55,0.9,SOFTGRN,GREEN,"local",GREEN,11); onodeC(3.85,2.55,0.9,SOFTPUR,PURPLE,"global",PURPLE,10.5)
connect(s,2.45,2.55,3.4,2.55,TEAL,2.5)
txt(s,1.6,3.15,2.9,0.4,[[Rn("cùng query → KÉO GẦN",11,TEAL,True)]],align=PP_ALIGN.CENTER)
onodeC(2.0,4.35,0.9,SOFTGRN,GREEN,"local",GREEN,11); onodeC(4.35,4.65,0.9,RGBColor(0xEE,0xE0,0xE0),RGBColor(0xCC,0x55,0x55),"q'",RGBColor(0xCC,0x55,0x55),12)
connect(s,2.45,4.4,3.9,4.6,RGBColor(0xCC,0x55,0x55),2.0)
txt(s,1.7,5.2,3.4,0.4,[[Rn("khác query → ĐẨY XA",11,RGBColor(0xCC,0x55,0x55),True)]],align=PP_ALIGN.CENTER)
formula(s,5.9,1.6,6.6,1.4,[
 "Cặp dương: (local, global) của CÙNG query",
 "Cặp âm: query KHÁC nhau",
 "L_cl = L_lg + L_gl + L_ll + L_gg"],title="4 contrastive loss")
bullets(s,5.9,3.2,6.6,2.0,[
 "Local dễ nhiễu; global ổn định.",
 "Kéo 2 view lại gần buộc local học đặc trưng BỀN, không phụ thuộc chi tiết bề mặt dễ nhiễu.",
 "→ Mô hình chống noise tốt hơn hẳn."],col=INK,bcol=RUST,size=13.5,gap=9)
example_panel(s,5.9,5.35,6.6,1.0,"Hiệu quả",[
 {'runs':[Rn("Khi thêm noise, LogCL giảm ÍT hơn nhiều so với RE-GCN/TiRGN.",12.5,INK)]},
])
note("Thành phần 3 (Contrastive): ép biểu diễn cùng query từ local & global giống nhau (dương), query khác thì khác (âm); 4 loss Llg/Lgl/Lll/Lgg. Chống noise vì global ổn định kéo local học đặc trưng bền.")

# ===================================================== 11 Ví dụ end-to-end
s=slide(); bg(s,LIGHT); header(s,"Phần 2 · LogCL","Ví dụ end-to-end: (China, Cooperate, ?, t_q)",BLUE)
rrect(s,0.85,1.7,3.75,4.55,SOFTGRN,GREEN,1.0)
txt(s,1.05,1.85,3.4,0.45,[[Rn("Local encoder",13.5,GREEN,True,f=HF)]])
txt(s,1.05,2.4,3.45,3.6,[
 {'runs':[Rn("Nhìn 7 snapshot gần.",12.5,INK,True)],'space_after':7},
 {'runs':[Rn("t_q−2: China vắng → α thấp.",12.5,INK)],'space_after':7},
 {'runs':[Rn("t_q−1: China “Make a visit” → α cao.",12.5,INK)],'space_after':7},
 {'runs':[Rn("GRU tổng hợp chuỗi → ",12.5,INK),Rn("h_tq",12.5,GREEN,True)],'space_after':0},
])
rrect(s,4.75,1.7,3.75,4.55,SOFTPUR,PURPLE,1.0)
txt(s,4.95,1.85,3.4,0.45,[[Rn("Global encoder",13.5,PURPLE,True,f=HF)]])
txt(s,4.95,2.4,3.45,3.6,[
 {'runs':[Rn("Lấy toàn lịch sử liên quan China + Cooperate.",12.5,INK)],'space_after':7},
 {'runs':[Rn("Subgraph: đối tác hợp tác lâu dài.",12.5,INK)],'space_after':7},
 {'runs':[Rn("Attention chọn entity lặp nhiều trong pattern Cooperate.",12.5,INK)],'space_after':7},
 {'runs':[Rn("→ ",12.5,INK),Rn("h_g",12.5,PURPLE,True)]},
])
rrect(s,8.65,1.7,3.8,4.55,SOFTGOLD,GOLD,1.0)
txt(s,8.85,1.85,3.45,0.45,[[Rn("Prediction",13.5,GOLD,True,f=HF)]])
txt(s,8.85,2.4,3.45,3.6,[
 {'runs':[Rn("Kết hợp theo λ (≈0.9, thiên local):",12.5,INK,True)],'space_after':6},
 {'runs':[Rn("ĥ = λ·h_tq + (1−λ)·h_g",12.5,INK,True,f=CF)],'space_after':10},
 {'runs':[Rn("Đưa vào ConvTransE →",12.5,INK)],'space_after':4},
 {'runs':[Rn("xác suất cho MỌI entity.",12.5,GOLD,True)]},
])
note("Ví dụ end-to-end (China, Cooperate): local — tq-2 China vắng α thấp, tq-1 'Make a visit' α cao → h_tq; global — subgraph đối tác lâu dài → h_g; prediction kết hợp λ≈0.9 thiên local → ConvTransE → xác suất mọi entity.")

# ===================================================== 12 Kết quả & ablation
s=slide(); bg(s,LIGHT); header(s,"Phần 2 · LogCL","Kết quả & ablation",BLUE)
rrect(s,0.85,1.75,5.6,2.0,SOFTBLUE,SOFTBLUE)
txt(s,1.1,1.95,5.1,0.45,[[Rn("Hiệu năng",14,BLUE,True,f=HF)]])
txt(s,1.1,2.5,5.15,1.1,[
 {'runs':[Rn("MRR cao hơn baseline tốt nhất (HisMatch) ",13.5,INK),Rn("4.9%–7.9%",14,BLUE,True),
          Rn(" trên 4 dataset; và giảm ÍT hơn nhiều khi thêm noise.",13.5,INK)]},
])
rrect(s,0.85,3.95,5.6,2.4,SOFTRUST,SOFTRUST)
txt(s,1.1,4.15,5.1,0.45,[[Rn("Vì sao bền với noise",14,RUST,True,f=HF)]])
txt(s,1.1,4.7,5.15,1.5,[
 {'runs':[Rn("Module contrastive ép local học đặc trưng bền theo global ổn định → ít “lung lay” trước nhiễu (khác hẳn RE-GCN/TiRGN).",13.5,INK)]},
])
rrect(s,6.7,1.75,5.75,4.6,WHITE,LINEC,1)
txt(s,6.95,1.95,5.3,0.45,[[Rn("Ablation (đóng góp từng phần)",14,INK,True,f=HF)]])
abl=[("Bỏ Global (LogCL-L)","− 4–5% MRR",PURPLE),
     ("Bỏ Entity-Aware Attention","− 8–10% MRR  (lớn nhất)",BLUE),
     ("Bỏ Contrastive (w/o-cl)","≈ ngang lúc thường, TỤT rõ khi có noise",RUST)]
yy=2.55
for t,v,c in abl:
    rrect(s,6.95,yy,5.25,1.05,LIGHT,LINEC,1); rect(s,6.95,yy,0.09,1.05,c)
    txt(s,7.2,yy+0.13,4.9,0.4,[[Rn(t,13.5,INK,True)]])
    txt(s,7.2,yy+0.55,4.9,0.4,[[Rn(v,13,c,True)]])
    yy+=1.2
note("Kết quả: MRR vượt HisMatch 4.9–7.9% trên 4 dataset, bền noise hơn. Ablation: bỏ global -4~5%; bỏ entity-aware attention -8~10% (lớn nhất); bỏ contrastive ít ảnh hưởng lúc thường nhưng tụt rõ khi có noise.")


# ===================================================== 13 DIVIDER P3 DEPTH
divider("3","DEPTH","Hợp nhất nhúng + đường: 1 động lực · 3 thành phần · end-to-end",GOLD)
note("Chuyển sang DEPTH — mô hình đề xuất, trình bày song song với LogCL: động lực, sơ đồ, 3 thành phần, ví dụ, ablation.")

# ===================================================== 14 Động lực
s=slide(); bg(s,LIGHT); header(s,"Phần 3 · DEPTH","Động lực: hai họ suy luận BÙ NHAU",GOLD)
rrect(s,0.85,1.75,5.6,4.0,SOFTBLUE,SOFTBLUE); rect(s,0.85,1.75,5.6,0.1,BLUE)
txt(s,1.1,1.95,5.1,0.45,[[Rn("Nhánh Nhúng (LogCL)",16,BLUE,True,f=HF)]])
for i,(l,good) in enumerate([("Điểm DÀY — phủ mọi thực thể",True),("Mạnh: thống kê, tương đồng, tần suất",True),
                             ("Yếu: thiên lệch độ phổ biến",False),("Transductive (cần entity-id)",False)]):
    txt(s,1.3,2.6+i*0.72,5.0,0.6,[[Rn(("✓  " if good else "✗  "),14,(BLUE if good else RGBColor(0xCC,0x55,0x55)),True),Rn(l,13.5,INK)]])
rrect(s,6.9,1.75,5.6,4.0,SOFTTEAL,SOFTTEAL); rect(s,6.9,1.75,5.6,0.1,TEAL)
txt(s,7.15,1.95,5.1,0.45,[[Rn("Nhánh Đường (CognTKE)",16,TEAL,True,f=HF)]])
for i,(l,good) in enumerate([("Điểm THƯA — chỉ thực thể có đường",True),("Mạnh: suy luận cấu trúc đa-hop, có căn cứ",True),
                             ("Yếu: = 0 khi không có đường",False),("Inductive (chỉ dùng quan hệ)",True)]):
    txt(s,7.35,2.6+i*0.72,5.0,0.6,[[Rn(("✓  " if good else "✗  "),14,(TEAL if good else RGBColor(0xCC,0x55,0x55)),True),Rn(l,13.5,INK)]])
rrect(s,0.85,5.95,11.6,0.85,SOFTGOLD,GOLD,1.0)
txt(s,1.1,6.05,11.1,0.65,[[Rn("Hai nhánh SAI ở các truy vấn KHÁC nhau (lỗi trực giao) → gộp lại thì bù trừ; điểm yếu “thiên lệch phổ biến” của nhúng được nhánh đường sửa.",13.5,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
note("Động lực DEPTH: nhúng (dày, thống kê, transductive, thiên lệch phổ biến) và đường (thưa, cấu trúc, inductive) bù nhau; lỗi trực giao nên gộp thì hợp hai tập đúng lớn hơn.")

# ===================================================== 15 SƠ ĐỒ kiến trúc DEPTH
s=slide(); bg(s,LIGHT); header(s,"Phần 3 · DEPTH","Sơ đồ kiến trúc tổng thể",GOLD)
# input
rrect(s,1.6,1.4,10.1,0.5,RGBColor(0xF1,0xEF,0xE8),RGBColor(0x5F,0x5E,0x5A),0.75)
txt(s,1.6,1.4,10.1,0.5,[[Rn("Truy vấn (s, r, ?, t)  ·  lịch sử K snapshot",12.5,INK,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
darrow(s,6.46,1.92,0.32,0.28)
# shared backbone bar
rrect(s,1.6,2.28,10.1,0.6,NAVY,NAVY,1.0)
txt(s,1.6,2.28,10.1,0.6,[[Rn("SHARED BACKBONE — RE-GCN + GRU · emb_rel DÙNG CHUNG  →  H ∈ ℝ^(N×d)",12.5,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
darrow(s,3.45,2.92,0.3,0.28); darrow(s,9.5,2.92,0.3,0.28)
# Head A
rrect(s,0.9,3.3,5.5,1.45,SOFTBLUE,BLUE,1.0)
txt(s,1.1,3.42,5.1,0.4,[[Rn("Head A · Embedding  (= LogCL)",13,BLUE,True,f=HF)]],align=PP_ALIGN.CENTER)
rrect(s,1.15,3.92,5.0,0.62,WHITE,BLUE,0.75)
txt(s,1.15,3.92,5.0,0.62,[[Rn("ConvTransE  →  S_embed ∈ ℝ^(B×N)  (DÀY)",11.5,BLUE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# Head B
rrect(s,6.95,3.3,5.5,1.45,SOFTTEAL,TEAL,1.0)
txt(s,7.15,3.42,5.1,0.4,[[Rn("Head B · Path  (mới, inline)",13,TEAL,True,f=HF)]],align=PP_ALIGN.CENTER)
rrect(s,7.2,3.86,5.0,0.4,WHITE,TEAL,0.75); txt(s,7.2,3.86,5.0,0.4,[[Rn("① boundary init tại subject s",11,TEAL,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
rrect(s,7.2,4.32,5.0,0.4,WHITE,TEAL,0.75); txt(s,7.2,4.32,5.0,0.4,[[Rn("② L=2 hop → S_path ∈ ℝ^(B×N) (THƯA)",11,TEAL,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# fusion
connect(s,3.6,4.75,5.6,4.98,GOLD,1.6); connect(s,9.7,4.75,7.5,4.98,GOLD,1.6)
rrect(s,4.35,4.98,4.6,0.6,SOFTGOLD,GOLD,1.0)
txt(s,4.35,4.98,4.6,0.6,[[Rn("FUSION:  S = S_embed + γ · z(S_path)",12.5,INK,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
darrow(s,6.46,5.6,0.3,0.26)
rrect(s,5.0,5.92,3.3,0.48,NAVY2,NAVY2); txt(s,5.0,5.92,3.3,0.48,[[Rn("softmax → ranking",12,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
darrow(s,6.46,6.42,0.3,0.24)
rrect(s,4.1,6.7,5.1,0.46,RGBColor(0xEA,0xF3,0xDE),RGBColor(0x3B,0x6D,0x11),0.75)
txt(s,4.1,6.7,5.1,0.46,[[Rn("L = L_object + L_contrastive + L_static  (end-to-end)",11.5,RGBColor(0x3B,0x6D,0x11),True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
note("Sơ đồ DEPTH (song song style LogCL): TKG → Shared Backbone (H[N×d], emb_rel chung) → Head A (ConvTransE→S_embed dày) ‖ Head B (boundary→L hop→S_path thưa) → Fusion S=S_embed+γz(S_path) → softmax→ranking; loss end-to-end L=L_object+L_cl+L_static.")

# ===================================================== 16 Thành phần 1 backbone+HeadA
s=slide(); bg(s,LIGHT); header(s,"Phần 3 · DEPTH · Thành phần 1","Backbone dùng chung + Head A (Embedding)",BLUE)
bullets(s,0.85,1.7,7.1,2.7,[
 "Cùng RE-GCN+GRU như LogCL tạo biểu diễn thực thể tiến hóa H ∈ ℝ^(N×d).",
 "Bảng quan hệ emb_rel ∈ ℝ^(2R×d) là “cầu nối” chung cho cả hai head.",
 "Head A = decoder ConvTransE của LogCL: ghép (ĥ_s, r) → chấm điểm mọi thực thể.",
 "Kết quả: S_embed — điểm DÀY cho toàn bộ 7128 thực thể."],col=INK,bcol=BLUE,size=13.5,gap=10)
formula(s,8.15,1.7,4.35,2.0,[
 "H = GRU-RGCN(G(t-7..t-1))",
 "S_embed[o] = ConvTransE(ĥ_s, r) · h_o",
 "                    ∈ ℝ^(B×N)"],title="Head A")
example_panel(s,8.15,3.95,4.35,2.4,"Ví dụ (USA, Sanction)",[
 {'runs':[Rn("S_embed: Russia 2.40 (#1), China 2.30, Iran 2.10 (#3), N.Korea 1.50.",12.5,INK)],'space_after':6},
 {'runs':[Rn("Iran đúng nhưng bị tụt do thiên lệch phổ biến.",12.5,GOLD,True)]},
])
note("Thành phần 1: backbone chung (H[N×d]) + emb_rel chung; Head A = ConvTransE của LogCL → S_embed dày. Ví dụ: Iran #3.")

# ===================================================== 17 Thành phần 2 Head B
s=slide(); bg(s,LIGHT); header(s,"Phần 3 · DEPTH · Thành phần 2","Head B — Path reasoner (NBFNet / RED-GNN)",TEAL)
formula(s,0.85,1.6,6.7,2.7,[
 "boundary:  h_v^(0) = 1[v=s] · W_b·emb_rel[r]",
 "",
 "lan truyền l=1..L:",
 " agg_v = (1/deg_v) Σ_(u→v) h_u^(l-1) ⊙ W_r·emb_rel[rel]",
 " h_v^(l) = ReLU( W_u·[agg_v ; h_v^(l-1)] ) + h_v^(0)",
 "",
 "điểm:  S_path[o] = w_out · h_o^(L)"],title="Lan truyền có điều kiện từ subject")
bullets(s,0.85,4.5,6.7,1.9,[
 "Chỉ node s “bật” nhãn ban đầu theo quan hệ truy vấn r.",
 "Mỗi hop truyền tin dọc cạnh, trọng số theo quan hệ.",
 "Chỉ dùng emb_rel → INDUCTIVE; rất nhẹ (dim=32, L=2)."],col=INK,bcol=TEAL,size=13,gap=7)
# mini graph
def onodeD(cx,cy,d,fill,line,label,tc,sz=11):
    shp(s,MSO_SHAPE.OVAL,cx-d/2,cy-d/2,d,d,fill,line,1.5)
    txt(s,cx-d/2,cy-d/2,d,d,[[Rn(label,sz,tc,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
S0=(8.6,2.6); M1=(10.3,2.05); M2=(10.3,3.15); A0=(12.0,2.6)
connect(s,*S0,*M1,TEAL,1.8); connect(s,*S0,*M2,TEAL,1.8); connect(s,*M1,*A0,TEAL,1.8); connect(s,*M2,*A0,TEAL,1.8)
onodeD(*S0,0.78,GOLD,GOLD,"USA",NAVY,10); onodeD(*M1,0.72,WHITE,TEAL,"Israel",TEAL,9)
onodeD(*M2,0.72,WHITE,TEAL,"IAEA",TEAL,9.5); onodeD(*A0,0.82,GOLD,GOLD,"Iran",NAVY,10)
rrect(s,7.75,3.75,4.75,2.65,SOFTGOLD,GOLD,1.0)
txt(s,7.95,3.9,4.4,0.4,[[Rn("S_path (thô) → z-norm",12.5,GOLD,True,f=HF)]])
sp=[("Iran","3.0 → z≈+2.6",TEAL),("N.Korea","1.0 → z≈+0.8",INK),("Russia","0 → z≈−0.1",INK),("China","0 → z≈−0.1",INK)]
for i,(nm,v,c) in enumerate(sp):
    txt(s,8.1,4.4+i*0.46,1.6,0.4,[[Rn(nm,12.5,c,c==TEAL)]])
    txt(s,9.55,4.4+i*0.46,2.8,0.4,[[Rn(v,12,c,c==TEAL,f=CF)]])
note("Thành phần 2 (Head B): công thức boundary 1[v=s]W_b r; layer agg=(1/deg)Σ h_u⊙W_r r, h=ReLU(W[agg;h])+boundary; score w_out·h^L. Inductive, nhẹ. Mini graph: USA→Israel/IAEA→Iran, hai đường → S_path[Iran]=3.0→z+2.6; Russia/China=0.")

# ===================================================== 18 Thành phần 3 Fusion
s=slide(); bg(s,LIGHT); header(s,"Phần 3 · DEPTH · Thành phần 3","Fusion γ & huấn luyện end-to-end",GOLD)
formula(s,0.85,1.55,11.6,1.1,[
 "S = S_embed + γ · z(S_path),   z(x)=(x−μ_row)/σ_row,   γ ∈ ℝ học được (init 0)"],title="Công thức hợp nhất")
cards=[("γ khởi tạo 0","Lúc đầu S = S_embed = LogCL (sàn an toàn). γ chỉ lớn lên nếu nhánh đường giảm loss.",BLUE),
       ("Chuẩn hóa z()","Đưa S_path về cùng thang với S_embed theo từng hàng, tránh một nhánh lấn át.",TEAL),
       ("End-to-end","Backbone+2 head+γ học chung; Head B nhận gradient → tự thích nghi bù lỗi nhúng.",GOLD)]
x=0.85;w=3.78;g=0.13
for i,(t,d,c) in enumerate(cards):
    cx=x+i*(w+g); rrect(s,cx,2.95,w,2.05,WHITE,LINEC,1); rect(s,cx,2.95,w,0.1,c)
    txt(s,cx+0.22,3.15,w-0.44,0.5,[[Rn(t,15,INK,True,f=HF)]]); txt(s,cx+0.22,3.75,w-0.44,1.15,[[Rn(d,12.5,MUTE)]])
formula(s,0.85,5.25,11.6,1.1,["L = L_object (NLL) + L_contrastive + L_static  — tối ưu toàn bộ mô hình cùng một lượt"],title="Hàm mất mát")
note("Thành phần 3 (Fusion): S=S_embed+γz(S_path); z-norm theo hàng; γ init 0 = sàn LogCL; end-to-end. Loss tổng.")

# ===================================================== 19 Ví dụ hợp nhất (payoff)
s=slide(); bg(s,LIGHT); header(s,"Phần 3 · DEPTH","Ví dụ end-to-end: hợp nhất đưa Iran lên #1",GOLD)
rows=[["Ứng viên","S_embed","z(S_path)","S = S_embed + 0.2·z","Hạng"],
      ["Russia","2.40","−0.1","2.38","2"],
      ["China","2.30","−0.1","2.28","3"],
      ["Iran  (đúng)","2.10","+2.6","2.62","1 ✓"],
      ["N.Korea","1.50","+0.8","1.66","4"]]
gt=s.shapes.add_table(5,5,Inches(1.4),Inches(1.75),Inches(10.5),Inches(3.0)).table
widths=[2.6,1.9,1.9,2.7,1.4]
for ci,wd in enumerate(widths): gt.columns[ci].width=Inches(wd)
for ri in range(5): gt.rows[ri].height=Inches(0.6)
for ri,row in enumerate(rows):
    iran=(ri==3)
    for ci,val in enumerate(row):
        cell=gt.cell(ri,ci); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.margin_left=Inches(0.1);cell.margin_right=Inches(0.05);cell.margin_top=Inches(0.02);cell.margin_bottom=Inches(0.02)
        if ri==0: fc,fillc,bold=WHITE,NAVY,True
        elif iran: fc,fillc,bold=NAVY,GOLD,True
        else: fc,fillc,bold=INK,WHITE,(ci==0)
        cell.fill.solid(); cell.fill.fore_color.rgb=fillc
        p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
        r=p.add_run(); r.text=val; r.font.size=Pt(13.5); r.font.bold=bold; r.font.color.rgb=fc
        r.font.name=CF if (ci in (1,2,3) and ri>0) else BF
txt(s,1.4,5.05,10.5,1.3,[
 {'runs':[Rn("Trước fusion: ",14,INK,True),Rn("Iran #3 (nhúng thiên lệch Russia). ",13.5,INK),
          Rn("Sau fusion: ",14,GOLD,True),Rn("đường bơm +0.2·2.6=+0.52 cho Iran → Iran #1. ",13.5,INK),
          Rn("Russia/China không có đường nên gần như không đổi.",13.5,MUTE)]},
])
note("Payoff: bảng S_embed | z(S_path) | fused. Iran 2.10→2.62 (#1); Russia 2.40→2.38. Lật #3→#1 nhờ nhánh đường.")

# ===================================================== 19a Động cơ (LogCL nền)
s=slide(); bg(s,LIGHT); header(s,"Phần 3 · DEPTH","Động cơ: lấy LogCL làm nền, vá điểm mù bằng path",GOLD)
rrect(s,0.85,1.7,6.5,4.55,SOFTBLUE,SOFTBLUE); rect(s,0.85,1.7,6.5,0.1,BLUE)
txt(s,1.1,1.9,6.0,0.45,[[Rn("LogCL — phương pháp NỀN",15,BLUE,True,f=HF)]])
txt(s,1.1,2.45,6.05,0.5,[[Rn("Mạnh: embedding tiến hóa + contrastive, đạt SOTA. Nhưng suy luận chủ yếu bằng độ tương đồng → còn 4 điểm mù:",13,INK)]])
for i,l in enumerate(["Thiên lệch độ phổ biến (hay đoán thực thể nổi tiếng)",
                      "Khó bắt quan hệ tổ hợp đa-hop (A→B→C)",
                      "Dự đoán KHÔNG có bằng chứng, khó giải thích",
                      "Transductive — học theo danh tính thực thể"]):
    txt(s,1.3,3.25+i*0.72,6.0,0.6,[[Rn("✗  ",13.5,RGBColor(0xCC,0x55,0x55),True),Rn(l,13,INK)]])
rrect(s,7.65,1.7,4.8,4.55,SOFTGOLD,GOLD,1.0);
txt(s,7.9,1.9,4.35,0.45,[[Rn("Ý tưởng",15,GOLD,True,f=HF)]])
txt(s,7.9,2.5,4.4,3.6,[
 {'runs':[Rn("Bổ sung cho LogCL đúng cái nó THIẾU: ",13.5,INK,True),
          Rn("khả năng suy luận theo ĐƯỜNG quan hệ đa-hop — thế mạnh của dòng CognTKE.",13.5,INK)],'space_after':12},
 {'runs':[Rn("→ Thiết kế một path-reasoning head nhẹ, ",13.5,INK),
          Rn("CẮM THẲNG vào LogCL",13.5,GOLD,True),
          Rn(" (không chạy CognTKE riêng).",13.5,INK)],'space_after':12},
 {'runs':[Rn("Path đóng vai “mô-đun bằng chứng cấu trúc” bổ trợ cho LogCL.",13,MUTE,it=True)]},
])
note("Động cơ: LogCL là nền mạnh nhưng có 4 điểm mù (thiên lệch phổ biến, khó đa-hop, không bằng chứng, transductive). Ý tưởng: cắm thêm một path-reasoning head (thế mạnh CognTKE) vào LogCL để vá đúng các điểm mù đó.")

# ===================================================== 19b Cơ chế gắn
s=slide(); bg(s,LIGHT); header(s,"Phần 3 · DEPTH","Cơ chế gắn: cắm mô-đun path vào LogCL",GOLD)
formula(s,0.85,1.6,7.05,4.6,[
 "# src/rrgcn.py  —  __init__()",
 "self.emb_rel   = Embedding(2R, d)     # DÙNG CHUNG",
 "self.path_head = PathHead(self.emb_rel, dim=32, L=2)",
 "self.gamma     = Parameter(zeros(1))  # init 0",
 "",
 "# get_loss() & predict()  — TRƯỚC softmax",
 "S_embed = decoder_ob(H, emb_rel, q)   # LogCL  [B,N]",
 "S_path  = self.path_head(s, r, G_hist)# path   [B,N]",
 "S       = S_embed + gamma*z_norm(S_path)  # FUSION",
 "P       = softmax(S)",
 "loss    = NLL(P,o) + L_cl + L_static",
 "loss.backward()      # END-TO-END (không freeze)"],title="Chỉ thêm 1 nhánh + 1 phép cộng vào code LogCL")
txt(s,8.1,1.6,4.4,0.45,[[Rn("Bốn thứ để “gắn”",14,TEAL,True,f=HF)]])
glue=[("① Điểm gắn","score-level, ngay trước softmax"),
      ("② Chất keo","cổng γ học được (init 0)"),
      ("③ Đồng thang","z-norm theo hàng cho cùng cỡ"),
      ("④ Cầu nối","dùng chung bảng emb_rel")]
for i,(t,d) in enumerate(glue):
    yy=2.15+i*1.0
    rrect(s,8.1,yy,4.4,0.85,WHITE,LINEC,1); rect(s,8.1,yy,0.09,0.85,TEAL)
    txt(s,8.35,yy+0.12,4.05,0.36,[[Rn(t,13,INK,True)]])
    txt(s,8.35,yy+0.46,4.05,0.34,[[Rn(d,11.5,MUTE)]])
note("Cơ chế gắn: score-level fusion S=S_embed+γ·z_norm(S_path) chèn trước softmax trong get_loss/predict; chất keo = cổng γ học được (init 0); z-norm đồng thang; cầu nối = dùng chung emb_rel; train end-to-end, không đụng lõi LogCL.")

# ===================================================== 19c Path bù gì + bằng chứng
s=slide(); bg(s,LIGHT); header(s,"Phần 3 · DEPTH","Path bù đúng cái LogCL thiếu — và đẩy điểm lên",GOLD)
txt(s,0.85,1.5,5.6,0.4,[[Rn("LogCL (nền) còn yếu",13.5,RGBColor(0xCC,0x55,0x55),True,f=HF)]])
txt(s,6.85,1.5,5.6,0.4,[[Rn("Path bổ sung",13.5,TEAL,True,f=HF)]])
pairs=[("Đoán theo độ phổ biến","Chỉ cần CÓ đường quan hệ nối tới → kéo đáp án hiếm-đúng lên"),
       ("Khó ghép đa-hop A→B→C","Lan truyền đúng theo cấu trúc nhiều bước"),
       ("Không có bằng chứng","Cho đường bằng chứng tường minh (USA→…→Iran)"),
       ("Phụ thuộc danh tính (transductive)","Chỉ dùng quan hệ (relation-only) → tổng quát hơn")]
y0=1.95; step=0.86
for i,(l,r) in enumerate(pairs):
    yy=y0+i*step
    rrect(s,0.85,yy,5.6,0.74,RGBColor(0xFA,0xEC,0xE7),RGBColor(0xCC,0x9A,0x90),1)
    txt(s,1.05,yy,5.2,0.74,[[Rn(l,12.5,INK)]],anchor=MSO_ANCHOR.MIDDLE)
    arrow(s,6.5,yy+0.22,0.3,0.3,TEAL)
    rrect(s,6.85,yy,5.6,0.74,SOFTTEAL,TEAL,1)
    txt(s,7.05,yy,5.2,0.74,[[Rn(r,12.5,INK)]],anchor=MSO_ANCHOR.MIDDLE)
rrect(s,0.85,5.55,11.6,1.05,NAVY,NAVY)
txt(s,1.15,5.7,11.1,0.8,[
 {'runs':[Rn("Bằng chứng: ",13.5,GOLD,True),
          Rn("tắt path (γ=0) = đúng LogCL 0.491  →  bật path (DEPTH) 0.527  =  +2.9 MRR / +3.0 Hits@1.  γ học ra DƯƠNG → mô hình tự nguyện dùng path.",13.5,ICE)]},
])
note("Path bù đúng 4 điểm mù của LogCL (bảng đối chiếu). Bằng chứng: ablation γ=0 (=LogCL 0.491) → bật path (DEPTH 0.527), +2.9 MRR/+3.0 Hits@1; γ hội tụ dương = path thật sự có ích.")

# ===================================================== 20 DIVIDER P4
divider("4","Hiệu năng","Kết quả trên ICEWS14 (time-aware filtered)",NAVY2)
note("Chuyển sang số liệu.")


# ===================================================== 22 Performance table
s=slide(); bg(s,LIGHT); header(s,"Phần 4 · Hiệu năng","Kết quả trên ICEWS14",GOLD)
rows=[["Mô hình","MRR","Hits@1","Hits@3","Hits@10"],
      ["CognTKE (paper)","46.06","36.49","51.11","64.49"],
      ["LogCL (paper)","48.87","37.76","54.71","70.26"],
      ["DEPTH (đề xuất)","52.72","41.89","58.61","73.21"],
      ["Δ so với LogCL","+3.85","+4.13","+3.90","+2.95"],
      ["Δ so với CognTKE","+6.66","+5.40","+7.50","+8.72"]]
gt=s.shapes.add_table(6,5,Inches(1.55),Inches(1.9),Inches(10.2),Inches(3.6)).table
gt.columns[0].width=Inches(3.4)
for ci in range(1,5): gt.columns[ci].width=Inches(1.7)
hr=[0.7,0.7,0.7,0.8,0.55,0.55]
for ri in range(6): gt.rows[ri].height=Inches(hr[ri])
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        cell=gt.cell(ri,ci); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.margin_left=Inches(0.1);cell.margin_right=Inches(0.06);cell.margin_top=Inches(0.02);cell.margin_bottom=Inches(0.02)
        if ri==0: fc,fillc,bold,sz=WHITE,NAVY,True,15
        elif ri==3: fc,fillc,bold,sz=NAVY,GOLD,True,15
        elif ri>=4: fc,fillc,bold,sz=ICE,NAVY2,True,12.5
        else: fc,fillc,bold,sz=INK,WHITE,(ci==0),14
        cell.fill.solid(); cell.fill.fore_color.rgb=fillc
        p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
        r=p.add_run(); r.text=val; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=fc
        r.font.name=CF if (ci>0 and 0<ri) else BF
txt(s,1.55,5.75,10.2,0.95,[
 {'runs':[Rn("DEPTH thắng CẢ 4 chỉ số. ",14,GOLD,True),
          Rn("Hai paper được tự tái lập trên cùng code/eval (LogCL Hits@10≈70.3=70.26; CognTKE=64.49) → cùng hệ quy chiếu.",13,MUTE)]},
])
note("Bảng kết quả ICEWS14, 4 chỉ số, 3 mô hình + 2 dòng delta. DEPTH thắng tất cả. Tái lập paper khớp → đáng tin.")

# ===================================================== 23 Conclusion
s=slide(); bg(s,NAVY); rect(s,0,0,13.333,0.18,GOLD)
txt(s,0.9,0.7,11.5,0.7,[[Rn("Kết luận",34,WHITE,True,f=HF)]])
pts=[("Tốt hơn","Thắng cả 4 chỉ số: vượt CognTKE +6.66 và LogCL +3.85 MRR.",GOLD),
     ("Gọn hơn","MỘT mô hình huấn luyện end-to-end, không cần hai mô hình rời.",TEAL),
     ("Tinh tế hơn","1 backbone chung · 2 head trực giao · fusion γ học được.",BLUE)]
y=1.8
for t,d,c in pts:
    rect(s,0.9,y+0.05,0.18,0.92,c); txt(s,1.25,y,3.4,0.5,[[Rn(t,20,GOLD if c==GOLD else c,True,f=HF)]])
    txt(s,1.25,y+0.5,11.0,0.6,[[Rn(d,15.5,ICE)]]); y+=1.2
rrect(s,0.9,5.6,11.5,1.05,NAVY2,NAVY2)
txt(s,1.2,5.78,11.0,0.8,[
 {'runs':[Rn("Lưu ý trung thực: ",13.5,GOLD,True),
          Rn("Head B inductive, nhưng Head A dùng embedding thực thể nên TỔNG THỂ vẫn transductive — đúng thiết lập chuẩn ICEWS14.",13.5,ICE)]},
])
note("Kết luận 3 ý + lưu ý trung thực về transductive.")

for _sl,_n in zip(prs.slides,NOTES): _sl.notes_slide.notes_text_frame.text=_n
import time as _t
_out=r"F:\New folder\CognTKE-main\CognTKE-main\DEPTH_slides.pptx"
try: prs.save(_out)
except PermissionError:
    _out=r"F:\New folder\CognTKE-main\CognTKE-main\DEPTH_slides_v3_%s.pptx"%_t.strftime("%H%M%S"); prs.save(_out)
print("OUT:",_out,"| slides:",len(prs.slides._sldIdLst),"| notes:",len(NOTES))
